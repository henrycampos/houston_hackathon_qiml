# -*- coding: utf-8 -*-
"""
Created on Tue Jan 22 08:31:42 2019

@author: N.Srisutthiyakorn
"""

# -*- coding: utf-8 -*-
"""
Created on Mon Jan 21 23:06:05 2019

@author: N.Srisutthiyakorn
"""
# =============================================================================
# ReadDataset.py
from mpl_toolkits.mplot3d import Axes3D
import matplotlib.pyplot as plt
from matplotlib import cm
from matplotlib.ticker import LinearLocator, FormatStrFormatter
import numpy as np
import scipy.interpolate
import os
import nserp as rp

# Data information
# data size 150x200x200 with cell dimensions 25x25x1m
# reservoir size 3750m (E-W) x 5000m (N-S) x 200m
# Three layer thickness 80x40x80 thick
# In reshape (Z, Y, X) or (Depth, N-S, E-W)

# =============================================================================
# Functions
def is_float(string):
    """ True if given string is float else False"""
    try:
        return float(string)
    except ValueError:
        return False
    
def loadDat(string, datShape):
    """ Load .dat data and convert them to numpy """
    # Shape of the data
    nz = datShape[0]
    ny = datShape[1]
    nx = datShape[2]
    
    # Open .dat file into numpy
    data = []
    with open(string, 'r') as f:
        d = f.readlines()
        for i in d:
            k = i.rstrip().split(",")
            data.append([float(i) if is_float(i) else i for i in k]) 
    
    # Remove any comments before since we know that the data start at row 3.
    data    = data[3:]
    data    = np.array(data, dtype='O')
    data    = np.reshape(data, (nz, ny, nx))
    
    return data
    
def plotSlice(data, slice_interval, fig_mainDir, dataName):
    """ Plot multiple slices from the 3-D volumes and save them into specified folder """
    
    # Get the size of the data
    nz = data.shape[0]
    
    # Create the folder to store files
    fig_subDir = os.path.join(fig_mainDir, dataName)
    if not os.path.exists(fig_subDir):
        os.makedirs(fig_subDir)
        
    # Plot each slice
    listDisplay = np.arange(0,nz,slice_interval)
    for iDisplay in listDisplay:
        fig = plt.figure(figsize = (4,5), facecolor = "white")
        #fig = plt.figure(facecolor = "white")
        data_slice = data[iDisplay,:,:].astype(float)
        plt.imshow(data_slice)
        plt.title(str(iDisplay))
        plt.colorbar()
        
        fig_fn = dataName + "_" + str(iDisplay).zfill(3) + ".png"
        fig_fn = os.path.join(fig_subDir, fig_fn)
        plt.savefig(fig_fn)    
    
def createModels(model):
    # separate the model into different layers
    model_layer1 = model[:80,:,:]
    model_layer2 = model[80:120,:,:]
    model_layer3 = model[120:,:,:]
    
    # model2 flip - retrogradational model
    model2 = np.flip(model, 0)
    
    # shuffle layers into different cases - shuffle between depositional sequence
    model3 = np.concatenate((model_layer1, model_layer2, model_layer1, model_layer2), axis = 0)
    model4 = np.concatenate((model_layer2, model_layer3, model_layer2, model_layer3), axis = 0)
    model5 = np.concatenate((np.flip(model_layer2,0), np.flip(model_layer1,0), np.flip(model_layer2,0), np.flip(model_layer1,0)), axis = 0)
    model6 = np.concatenate((np.flip(model_layer3,0), np.flip(model_layer2,0), np.flip(model_layer3,0), np.flip(model_layer2,0)), axis = 0)
        
    return model, model2, model3, model4, model5, model6
    
def calcVpVs(G, K, Rho):
    """Calculcate Vp Vs in m/s from G(GPa), K(GPa), and Rho(g/cc)"""
    # Unit conversion to Pa for moduli and Rho in kg/m3
    K = K * np.power(10,9)
    G = G * np.power(10,9)
    Rho = Rho * 1000

    Vp = np.sqrt((K + (4 * G / 3 )) / Rho)
    Vs = np.sqrt(G / Rho)
    
    return Vp, Vs

def calcAISI(Vp, Vs, Rho):
    AI = Rho*Vp
    SI = Rho*Vs
    
    return AI, SI

def convolve(tr, w):
    return np.convolve(tr, w, mode='same')



# =============================================================================        
# Reshape the data
nx = 150 # E-W
ny = 200 # N-S
nz = 200 # Depth
datShape = [nz, ny, nx]

# Main Data Directory
model_path  = r'C:\Users\N.Srisutthiyakorn\OneDrive - Shell\Documents\GitHub\ShellHackathon\SampleInputs\4QC'

# Data Loading for Porosity 
model_fn        = os.path.join(model_path, 'porosity.dat') 
model_porosity  = loadDat(model_fn, datShape)

# Data Loading for Facies
model_fn        = os.path.join(model_path, 'facies.dat') 
model_facies    = loadDat(model_fn, datShape)

# Data Loading for Vp 
model_fn        = os.path.join(model_path, 'Pvelocity.dat') 
model_Vp        = loadDat(model_fn, datShape)

# Data Loading for Vs 
model_fn        = os.path.join(model_path, 'Svelocity.dat') 
model_Vs        = loadDat(model_fn, datShape)

# Data Loading for Density
model_fn        = os.path.join(model_path, 'density.dat') 
model_den       = loadDat(model_fn, datShape)

# Data Loading for Acoustic Impedance 
model_fn        = os.path.join(model_path, 'acoustic_impedance.dat') 
model_AI        = loadDat(model_fn, datShape)

# Data Loading for Shear Impedance
model_fn        = os.path.join(model_path, 'Swave_impedance.dat') 
model_SI        = loadDat(model_fn, datShape)

# Data Loading for Elastic Impedance
model_fn        = os.path.join(model_path, 'Swave_impedance.dat') 
model_SI        = loadDat(model_fn, datShape)

# Data Loading for Poisson Ratio
model_fn        = os.path.join(model_path, 'Poisson_ratio.dat') 
model_PR        = loadDat(model_fn, datShape)

# QC the calculation - pass
test0 = model_AI[0,0,0] - model_den[0,0,0] * model_Vp[0,0,0] 
test1 = model_SI[0,0,0] - model_den[0,0,0] * model_Vs[0,0,0] 

# Test data save/load
model_path      = r'C:\Users\N.Srisutthiyakorn\OneDrive - Shell\Documents\GitHub\ShellHackathon\SampleInputs\IntermediateSteps'
model_fn_out    = os.path.join(model_path, 'density_output_test.npy') 
np.save(model_fn_out, model_den)
model_den_qc = np.load(model_fn_out)

# =============================================================================
# Test wavelet
from bruges.filters import ricker
delta_t = 0.001 # sample rate of the data
wavelet, tw = ricker(0.128, delta_t, 25.0, return_t=True)
plt.plot(wavelet,'.')

# =============================================================================
# Test synthetic models
from bruges.reflection import reflectivity

theta = np.arange(46)
model_Vp_trace = model_Vp[:,0,0]
model_Vs_trace = model_Vs[:,0,0]
model_den_trace = model_den[:,0,0]

rc_theta = reflectivity(model_Vp_trace, model_Vs_trace, model_den_trace, theta).T

plt.imshow(rc_theta.real, aspect='auto')

def convolve(tr, w):
    return np.convolve(tr, w, mode='same')

s = np.apply_along_axis(convolve, 0, rc_theta.real, wavelet)
plt.figure(figsize=(6, 10))
plt.imshow(s, cmap="RdBu", aspect='auto', clim=[-0.25, 0.25], extent=[theta[0], theta[-1], t[-1], t[0]])
plt.ylim(2.0, 1.2)
plt.colorbar()
plt.show()

fig, axs = plt.subplots(figsize=(7, 10),
                        ncols=3,
                        sharey=True,
                        gridspec_kw=dict(width_ratios=[1.5, 1, 1 ]),
                       )

# Plot synthetic gather.
ax = axs[0]
ax.imshow(s, cmap="seismic", aspect='auto', clim=(-0.35,0.35), extent=[0,60, t[-1], t[0]])
gain = 10
for i, tr in enumerate(s.T):
    if i % 2 == 1:
        axs[0].plot(gain*(tr)+i, t[:-1], 'k', alpha=0.5)
ax.set_xlim(0,45)
ax.set_ylim(2.0, 1.2)
ax.set_xlabel('two-way-time (s)')
ax.set_title('incidence angle ($\degree$)')

# Plot impedance log.
ax = axs[1]
ax.plot(model_Vp_trace*model_den_trace, t,  lw=1.0)
ax.set_xlim(np.percentile(imp,5)*0.8, np.percentile(imp,95)*1.2)
ax.grid(c='k', alpha=0.25)
ax.set_yticks([])
ax.set_title('impedance')

## Plot colour-filled GR.
#ax = axs[2]
#ax.plot(gr_t, t,  c='k',lw=1.0)
#ax.fill_betweenx(t, gr_t, 0, color='lightgrey')
#ax.fill_betweenx(t, gr_t, 100, color='khaki')
#ax.grid(c='k', alpha=0.25)
#ax.set_xlim(20,100)
#ax.set_yticks([])
#ax.set_xticks([25,50,75,100])
#ax.grid(lw=0.5)
#ax.set_title('gamma ray (API)')

plt.show()
## Data Loading for Top
#model_fn        = os.path.join(model_path, 'Top\\top.dat') 
#model_top       = loadDat(model_fn, [200, 150])


# =============================================================================
# Plot the data
fig_mainDir = r'C:\Users\N.Srisutthiyakorn\OneDrive - Shell\Documents\GitHub\ShellHackathon\Figures'
slice_interval = 10

plotSlice(model_porosity, slice_interval, fig_mainDir, "Porosity")
plotSlice(model_facies, slice_interval, fig_mainDir, "Facies")
plotSlice(model_top, slice_interval, fig_mainDir, "Top")
    


# =============================================================================
# Creating new geometry
model_porosity1, model_porosity2, model_porosity3, model_porosity4, model_porosity5, model_porosity6 = createModels(model_porosity)


# =============================================================================
# Facies/Lithology Model
# 4 Facies in total (floodplain, point bar, channel, boundary)

# Mineral Fraction of each facies in order(clay, quartz, feldspar, rock fragments)
facies_minFract = np.array([[0.85, 0.15, 0, 0], [0 ,0.7,0.2,0.1], [0.0, 0.65, 0.2, 0.15], [0.9, 0.1, 0.0, 0.0]])

# Density and elastic properties of the mineral in the order (clay, quartz, feldspar, rock fragments) (Mavko et al, 2009)
min_rhob_gcc    = np.array([2.5, 2.65, 2.63, 2.7])
min_K_GPa       = np.array([21, 36.6, 75.6, 80])
min_G_GPa       = np.array([9, 44, 25.6, 20])
## Mineral elastic properties in range - future cases
#min_K_GPa_clayRange = [21, 25]
#min_G_GPa_clayRange = [7, 9]
#
#min_K_GPa_QtzRange = [36.5, 37.9]
#min_G_GPa_QtzRange = [44.3, 45.6]


# Mix rocks at different constituent - Voigt and Ruess average - use wider range than H-S bound
rock_K_GPa_ub = np.matmul(facies_minFract, min_K_GPa)
rock_K_GPa_lb = 1/(np.matmul(facies_minFract, 1/min_K_GPa))

rock_G_GPa_ub = np.matmul(facies_minFract, min_G_GPa)
rock_G_GPa_lb = 1/(np.matmul(facies_minFract, 1/min_G_GPa))

rock_rhob_gcc   = np.matmul(facies_minFract, min_rhob_gcc)

rock_Vp_ub, rock_Vs_ub  = calcVpVs(rock_K_GPa_ub, rock_G_GPa_ub, rock_rhob_gcc)
rock_Vp_lb, rock_Vs_lb  = calcVpVs(rock_K_GPa_lb, rock_G_GPa_lb, rock_rhob_gcc)


# =============================================================================
# Fluid Model
# Pressure (MPa) and Temperature(c) - assuming the constant for now
wd_m    = 500
dbml_m  = 2500 

P = (wd_m * 0.01) +  0.0135 * dbml_m
T = (0.023 * dbml_m) + 20.3

# Brine
Sa = 20000 # Brine Salinity (ppm)
rhob, Kb, vpb = rp.flagBrine(Sa, T, P)

# Gas
G = 0.70 # Gas Gravity
rhog, Kg, vpg = rp.flagGas(G, T, P) 

# Oil 
API = 25; GOR = 200;
Pbp, rhoo, Ko, vpo = rp.flagOil(G, T, P, API, GOR)
Ko[P < Pbp] = np.NaN; rhoo[P < Pbp] = np.NaN


# =============================================================================
# Fluid Mixing Model - keep this fixed for now.
# Case 1: Low Gas Saturation (Residual Gas)
Sg = 0.1; So = 0.0
Kreuss_lg, Kvoigt_lg, rhoeff_lg = rp.mixer(So, Sg, rhob, rhoo, rhog, Kb, Ko, Kg)

# Case 2: High Gas Saturation 
Sg = 0.9; So = 0.0
Kreuss_hg, Kvoigt_hg, rhoeff_hg = rp.mixer(So, Sg, rhob, rhoo, rhog, Kb, Ko, Kg)

# Case 3: Oil 
Sg = 0.0; So = 0.8
Kreuss_o, Kvoigt_o, rhoeff_o = rp.mixer(So, Sg, rhob, rhoo, rhog, Kb, Ko, Kg)

# Case 4: Wet
Sg = 0.0; So = 0.0
Kreuss_w, Kvoigt_w, rhoeff_w = rp.mixer(So, Sg, rhob, rhoo, rhog, Kb, Ko, Kg)


# =============================================================================
# Facies Model (Rock + Fluid) - Gassmann Substitution - must randomize the Vp Vs already
# phi would be varying according to the model
# Sw would varying in different fluid cases
# Question - K0 or Kdry - why not assuming the constituents rather - check
# rock_Vp has to be draw randomly and it also depends on facies. Need to loop to do the fluid substitution.
# Would be faster to have it in the vector? and plug them back again.
# rhoeff_lg, Kreuss_lg is varying from different cases

vec_porosity    = np.concatenate(np.concatenate(model_porosity))
vec_facies      = np.concatenate(np.concatenate(model_facies))
vec_facies      = vec_facies.astype(int)

nFacies         = np.unique(vec_facies)

# Fluid substitution needs to be done for each facies individually
count = 0
for iFacies in np.arange(4):
    print(iFacies)
    idx = vec_facies == iFacies
    count += np.count_nonzero(idx) 
    
print(count)
    
Vp_Case1, Vs_Case1, Rho_Case1, KMin_Case1 = rp.gassmnv(rock_Vp_ub, rock_Vs_ub, rock_rhob_gcc*1000, rhob*1000, Kb*1e+09, rhoeff_lg*1000, Kreuss_lg*1e+09, 37e+09, phi)






# =============================================================================
# AI and SI volumes - simply multiplication



# =============================================================================
# Run the Powerpoint Output
from pptx import Presentation
from pptx.util import Inches

# Powerpoint File Location
pp_path     = r'C:\Users\N.Srisutthiyakorn\OneDrive - Shell\Documents\GitHub\ShellHackathon'
pp_fnInput  = os.path.join(pp_path, 'ShellTemplate.pptx') # Generate this file from any powerpoint with Shell Template
pp_fnOutput = os.path.join(pp_path, 'Dataset.pptx')

# Open the powerpoint file
prs = Presentation(pp_fnInput)
blank_slide_layout = prs.slide_layouts[6]

# Figure Files Location (Main Folder)
fig_mainDir = r'C:\Users\N.Srisutthiyakorn\OneDrive - Shell\Documents\GitHub\ShellHackathon\Figures'

# Determining Figures Location in the slides
nRow        = 5
nCol        = 4
space_Row   = 1.2 # inch
space_Col   = 1.5
left_corner = 7 # inch
top_corner  = 0.8 
image_height    = 1.5 # inch
left_location   = np.arange(0, nCol*space_Col, space_Col) + left_corner
top_location    = np.arange(0, nRow*space_Row, space_Row) + top_corner
left_location, top_location = np.meshgrid(left_location, top_location)
left_location   = np.concatenate(left_location)
top_location    = np.concatenate(top_location)

# List Folders within the folder, one folder per slides
listFolder  = os.listdir(fig_mainDir)
for iFolder in listFolder:
    fig_subDir = os.path.join(fig_mainDir, iFolder)
    #print(fig_subDir)
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.title
    title.text = "Model - " + iFolder 
    
    # List Figures within the folder
    listFigure = os.listdir(fig_subDir)
    for count, iFile in enumerate(listFigure):
        fig_fn = os.path.join(fig_subDir, iFile)
        #print(fig_fn)
        pic = slide.shapes.add_picture(fig_fn, Inches(left_location[count]), Inches(top_location[count]), height = Inches(image_height))

prs.save(pp_fnOutput)
    
    



